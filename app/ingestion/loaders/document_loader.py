"""Unified multi-format document loader."""
from pathlib import Path
from typing import Dict, Any, List
from app.ingestion.loaders.base import BaseLoader, LoadedDocument

class DocumentLoader(BaseLoader):
    """Loads PDF, DOCX, TXT, MD, PPTX and raw text documents."""

    def load(self, file_path: str | Path) -> LoadedDocument:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {path}")

        ext = path.suffix.lower()
        content = ""
        meta = {
            "source": str(path),
            "filename": path.name,
            "extension": ext,
            "file_size": path.stat().st_size
        }

        if ext in (".txt", ".md", ".markdown", ".csv", ".json", ".log", ".py"):
            content = path.read_text(encoding="utf-8", errors="ignore")
            
        elif ext == ".pdf":
            try:
                import pypdf
                reader = pypdf.PdfReader(str(path))
                pages_text = []
                for i, page in enumerate(reader.pages):
                    t = page.extract_text() or ""
                    pages_text.append(t)
                content = "\n\n".join(pages_text)
                meta["page_count"] = len(reader.pages)
            except Exception as e:
                content = f"[Error reading PDF: {e}]"
                
        elif ext in (".docx", ".doc"):
            try:
                import docx
                doc = docx.Document(str(path))
                content = "\n\n".join([p.text for p in doc.paragraphs if p.text.strip()])
            except Exception as e:
                content = f"[Error reading DOCX: {e}]"
                
        elif ext in (".pptx", ".ppt"):
            try:
                from pptx import Presentation
                prs = Presentation(str(path))
                slide_texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            slide_texts.append(shape.text)
                content = "\n\n".join(slide_texts)
            except Exception as e:
                content = f"[Error reading PPTX: {e}]"
        else:
            # Attempt plain text read
            content = path.read_text(encoding="utf-8", errors="ignore")

        return LoadedDocument(content=content, metadata=meta, source_path=str(path))

    def load_directory(self, dir_path: str | Path, recursive: bool = True) -> List[LoadedDocument]:
        root = Path(dir_path)
        if not root.exists() or not root.is_dir():
            return []
            
        pattern = "**/*.*" if recursive else "*.*"
        docs = []
        for p in root.glob(pattern):
            if p.is_file() and not p.name.startswith("."):
                try:
                    docs.append(self.load(p))
                except Exception:
                    pass
        return docs

document_loader = DocumentLoader()
